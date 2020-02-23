from pptx import Presentation
import glob
import string
for eachfile in glob.glob("*.pptx"):
    filea = eachfile.replace(".pptx","")    
    file = open(filea+".sng","w", encoding="utf-8") 
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:         
            if hasattr(shape, "text"):
                if shape.text.isdigit():
                    print("Fool")
                else:   
                    file.write(shape.text)
                    file.write("\n---\n")
    file.close()
